import customtkinter as ctk
import os
import json
import shutil
import hashlib
import threading
import requests
import re
import docx
import openpyxl
from pptx import Presentation
from tkinter import messagebox, filedialog, simpledialog
from datetime import datetime
import pdfplumber
from mutagen.easyid3 import EasyID3
from mutagen.mp4 import MP4

# ==================== CONFIGURATION ====================
CONFIG_FILE = "ged_enterprise_config.json"
INDEX_FILE = "ged_file_index.json"
API_KEY = "api-key"
DEEPSEEK_API_URL = "https://api.deepseek.com/v1/chat/completions"

# ==================== CLASSES UTILITAIRES ====================
class ConfigManager:
    """Gestionnaire de configuration centralisé"""
    @staticmethod
    def load_config():
        if not os.path.exists(CONFIG_FILE):
            default_config = {
                "typology": {
                    "JURIDIQUE": ["Baux", "Actes"],
                    "TECHNIQUE": ["Diagnostics", "Visites_Video"],
                    "COMPTABILITE": ["Factures", "Audios_Etats_Lieux"],
                    "ADMINISTRATIF": ["Assurances", "Courriers", "Identité"]
                },
                "auto_delete": False,
                "last_destination": os.path.expanduser("~"),
                "api_active": True,
                "auto_create_categories": True  # Nouvelle option
            }
            ConfigManager.save_config(default_config)
            return default_config
        
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            print(f"Erreur chargement config: {e}")
            return {}

    @staticmethod
    def save_config(data):
        try:
            with open(CONFIG_FILE, "w", encoding="utf-8") as f:
                json.dump(data, f, indent=4, ensure_ascii=False)
        except Exception as e:
            print(f"Erreur sauvegarde config: {e}")

    @staticmethod
    def load_index():
        if not os.path.exists(INDEX_FILE):
            return {}
        try:
            with open(INDEX_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            print(f"Erreur chargement index: {e}")
            return {}

    @staticmethod
    def save_index(data):
        try:
            with open(INDEX_FILE, "w", encoding="utf-8") as f:
                json.dump(data, f, indent=4)
        except Exception as e:
            print(f"Erreur sauvegarde index: {e}")

class DuplicateManager:
    """Gestionnaire de détection de doublons par empreinte SHA-256"""
    @staticmethod
    def get_file_hash(filepath):
        """Calcule l'empreinte SHA-256 d'un fichier"""
        sha256_hash = hashlib.sha256()
        try:
            with open(filepath, "rb") as f:
                for byte_block in iter(lambda: f.read(4096), b""):
                    sha256_hash.update(byte_block)
            return sha256_hash.hexdigest()
        except Exception as e:
            print(f"Erreur calcul hash: {e}")
            return None

    @staticmethod
    def is_duplicate(file_hash, index_data):
        """Vérifie si un fichier existe déjà dans l'index"""
        return file_hash in index_data

class MetadataManager:
    """Gestionnaire des métadonnées pour fichiers audio/vidéo"""
    @staticmethod
    def tag_file(filepath, category, subcategory):
        """Injecte des métadonnées dans les fichiers"""
        ext = os.path.splitext(filepath)[1].lower()
        try:
            if ext == ".mp3":
                try:
                    audio = EasyID3(filepath)
                except:
                    audio = EasyID3()
                    audio.save(filepath)
                    audio = EasyID3(filepath)
                
                audio['genre'] = category
                audio['album'] = subcategory
                audio['artist'] = "MALKOGED AI"
                audio.save()
                
            elif ext in [".mp4", ".m4a", ".m4v"]:
                try:
                    video = MP4(filepath)
                except:
                    video = MP4()
                
                video["\xa9gen"] = category  # Tag Genre
                video["\xa9alb"] = subcategory  # Tag Album/Projet
                video["\xa9art"] = "MALKOGED AI"
                video.save()
                
        except Exception as e:
            print(f"Erreur tagging {filepath}: {e}")

class ClassificationEngine:
    """Moteur de classification IA DeepSeek avec création automatique de catégories"""
    def __init__(self):
        self.config = ConfigManager.load_config()
        self.typology = self.config.get("typology", {})
        self.api_available = self.config.get("api_active", True) and API_KEY and API_KEY != "TA_CLE_API_ICI"
        self.auto_create_categories = self.config.get("auto_create_categories", True)

    def reload_typology(self):
        """Recharge la typologie depuis le fichier de configuration"""
        self.config = ConfigManager.load_config()
        self.typology = self.config.get("typology", {})
        self.auto_create_categories = self.config.get("auto_create_categories", True)
        return self.typology

    def extract_text_from_pdf(self, filepath):
        """Extrait le texte d'un PDF"""
        text = ""
        try:
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages[:50]:  # Limité à 50 pages pour performance
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
        except Exception as e:
            print(f"Erreur extraction PDF {filepath}: {e}")
        return text

    def analyze_filename(self, filename):
        """Analyse le nom de fichier pour déterminer la catégorie"""
        filename_lower = filename.lower()
        
        # Règles de classification basées sur le nom
        rules = {
            "JURIDIQUE": ["bail", "acte", "contrat", "legal", "juridique"],
            "TECHNIQUE": ["diagnostic", "technique", "plan", "devis", "video", "photo"],
            "COMPTABILITE": ["facture", "compte", "bancaire", "impôt", "fiscal"],
            "ADMINISTRATIF": ["assurance", "courrier", "identite", "administratif"]
        }
        
        for category, keywords in rules.items():
            if any(keyword in filename_lower for keyword in keywords):
                return category
        
        return None  # Retourne None si aucune catégorie ne correspond
    
    def extract_text(self, filepath):
        """Extrait le texte de différents types de fichiers"""
        ext = os.path.splitext(filepath)[1].lower()
        text = ""
        try:
            if ext == ".pdf":
                text = self.extract_text_from_pdf(filepath)
            elif ext == ".docx":
                doc = docx.Document(filepath)
                text = "\n".join([para.text for para in doc.paragraphs])
            elif ext == ".xlsx":
                wb = openpyxl.load_workbook(filepath, read_only=True)
                # On lit les premières lignes de chaque feuille pour le contexte
                for sheet in wb.worksheets[:2]:
                    for row in sheet.iter_rows(max_row=20, values_only=True):
                        text += " ".join([str(cell) for cell in row if cell]) + "\n"
            elif ext == ".pptx":
                prs = Presentation(filepath)
                for slide in prs.slides[:5]:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
        except Exception as e:
            print(f"Erreur d'extraction sur {ext}: {e}")
        return text
    
    def call_deepseek_api(self, prompt_text):
        """Appelle l'API DeepSeek pour classification"""
        try:
            headers = {
                "Authorization": f"Bearer {API_KEY}",
                "Content-Type": "application/json"
            }
            
            payload = {
                "model": "deepseek-chat",
                "messages": [
                    {"role": "system", "content": "Tu es un assistant spécialisé dans la classification et l'indexation documentaires."},
                    {"role": "user", "content": prompt_text}
                ],
                "temperature": 0.1,
                "max_tokens": 500
            }
            
            response = requests.post(DEEPSEEK_API_URL, headers=headers, json=payload, timeout=30)
            response.raise_for_status()
            return response.json()["choices"][0]["message"]["content"]
            
        except Exception as e:
            print(f"Erreur API DeepSeek: {e}")
            return None

    def suggest_new_category(self, content_text, filename):
        """Demande à l'IA de suggérer une nouvelle catégorie et sous-catégorie"""
        try:
            prompt = f"""
            Analyse ce document pour créer une classification pertinente :
            
            Nom du fichier: {filename}
            
            Contenu (extrait):
            --- {content_text[:1500]} ---
            
            Tu es un expert en gestion documentaire et en classification documentaire.
            
            1. Analyse le document pour comprendre sa nature
            2. Propose une catégorie principale pertinente en fonction du contenu
            3. Propose une sous-catégorie spécifique
            
            Règles importantes :
            - Les catégories doivent être en MAJUSCULES
            - Les sous-catégories doivent être descriptives
            - Utilise un langage professionnel
            - Sois précis et concis
            
            Réponds UNIQUEMENT au format JSON :
            {{
                "category": "NOM_CATEGORIE_EN_MAJUSCULES",
                "subcategory": "Nom_Sous_Catégorie_Descriptif",
                "reason": "Brève explication du choix. Une ou deux phrases maximum."
            }}
            """
            
            result_text = self.call_deepseek_api(prompt)
            
            if result_text:
                clean_text = result_text.strip()
                if "```json" in clean_text:
                    clean_text = clean_text.split("```json")[1].split("```")[0].strip()
                elif "```" in clean_text:
                    clean_text = clean_text.split("```")[1].strip()
                
                json_match = re.search(r'\{.*\}', clean_text, re.DOTALL)
                if json_match:
                    try:
                        ai_suggestion = json.loads(json_match.group())
                        return ai_suggestion
                    except:
                        print("Erreur parsing JSON pour suggestion de catégorie")
        
        except Exception as e:
            print(f"Erreur suggestion catégorie: {e}")
        
        return None

    def auto_classify_with_creation(self, content_text, filename, existing_typology):
        """Classification avec création automatique de catégories"""
        # Essaie d'abord de trouver une catégorie existante
        filename_category = self.analyze_filename(filename)
        
        if filename_category and filename_category in existing_typology:
            # Cherche des mots-clés dans le contenu pour la sous-catégorie
            content_lower = content_text.lower()
            suggested_sub = self.suggest_subcategory_from_content(content_lower, filename_category, existing_typology)
            
            return {
                "category": filename_category,
                "subcategory": suggested_sub or existing_typology[filename_category][0] if existing_typology[filename_category] else "Divers",
                "created_new": False
            }
        
        # Si aucune catégorie existante ne correspond, crée une nouvelle
        if self.auto_create_categories:
            ai_suggestion = self.suggest_new_category(content_text, filename)
            
            if ai_suggestion:
                new_category = ai_suggestion.get("category", "AUTRE")
                new_subcategory = ai_suggestion.get("subcategory", "Divers")
                
                # Nettoyer le nom de catégorie
                new_category = new_category.strip().upper()
                new_subcategory = new_subcategory.strip()
                
                return {
                    "category": new_category,
                    "subcategory": new_subcategory,
                    "created_new": True,
                    "reason": ai_suggestion.get("reason", "")
                }
        
        # Fallback
        return {
            "category": "GENERAL",
            "subcategory": "Divers",
            "created_new": True,
            "reason": "Catégorie par défaut"
        }

    def suggest_subcategory_from_content(self, content_text, category, existing_typology):
        """Suggère une sous-catégorie basée sur le contenu"""
        if not content_text or len(content_text) < 50:
            return None
        
        # Règles de sous-catégories par catégorie
        rules = {
            "JURIDIQUE": {
                "bail": "Baux",
                "contrat": "Contrats",
                "acte": "Actes",
                "procès": "Contentieux",
                "tribunal": "Contentieux"
            },
            "TECHNIQUE": {
                "diagnostic": "Diagnostics",
                "devis": "Devis",
                "plan": "Plans",
                "photo": "Photos",
                "video": "Vidéos",
                "visite": "Visites"
            },
            "COMPTABILITE": {
                "facture": "Factures",
                "relevé": "Relevés",
                "impôt": "Impôts",
                "taxe": "Impôts",
                "bancaire": "Relevés_Bancaires"
            },
            "ADMINISTRATIF": {
                "assurance": "Assurances",
                "courrier": "Courriers",
                "identité": "Identité",
                "permis": "Permis",
                "autorisation": "Autorisations"
            }
        }
        
        if category in rules:
            for keyword, subcategory in rules[category].items():
                if keyword in content_text:
                    return subcategory
        
        return None

    def analyze_document(self, filepath):
        """Analyse un document et retourne sa classification avec création automatique de catégories si besoin"""
        filename = os.path.basename(filepath)
        
        # Classification initiale par nom de fichier
        predicted_category = self.analyze_filename(filename)
        predicted_sub = "Divers"
        created_new = False
        reason = ""
        
        # Extraction du contenu pour analyse approfondie
        content_text = ""
        supported_ext = ('.pdf', '.docx', '.xlsx', '.pptx')
        if filepath.lower().endswith(supported_ext):
            content_text = self.extract_text(filepath)
        
        # Si l'API est disponible et nous avons du contenu
        if self.api_available and len(content_text) > 10:
            try:
                # Classification intelligente avec création automatique
                classification_result = self.auto_classify_with_creation(
                    content_text, 
                    filename, 
                    self.typology
                )
                
                predicted_category = classification_result["category"]
                predicted_sub = classification_result["subcategory"]
                created_new = classification_result.get("created_new", False)
                reason = classification_result.get("reason", "")
                
                # Si une nouvelle catégorie a été créée, l'ajouter à la typologie
                if created_new and predicted_category not in self.typology:
                    self.typology[predicted_category] = [predicted_sub]
                    # Sauvegarder automatiquement la nouvelle typologie
                    self.config["typology"] = self.typology
                    ConfigManager.save_config(self.config)
                    print(f"Nouvelle catégorie créée: {predicted_category} > {predicted_sub}")
                
                # Si la catégorie existe mais pas la sous-catégorie, l'ajouter
                elif predicted_category in self.typology and predicted_sub not in self.typology[predicted_category]:
                    self.typology[predicted_category].append(predicted_sub)
                    self.config["typology"] = self.typology
                    ConfigManager.save_config(self.config)
                    print(f"Nouvelle sous-catégorie ajoutée: {predicted_category} > {predicted_sub}")
                    
            except Exception as e:
                print(f"Erreur analyse IA avec création: {e}")
                # Fallback sur la classification par nom
                if not predicted_category:
                    predicted_category = "GENERAL"
        
        # Si pas d'analyse IA possible, utiliser la classification par nom
        elif not predicted_category:
            predicted_category = "GENERAL"
        
        # Nommage standardisé
        doc_date = datetime.now().strftime("%Y%m%d")
        clean_filename = filename.replace(" ", "_").replace("(", "").replace(")", "")
        
        # Ajouter un marqueur si nouvelle catégorie créée
        status_prefix = "🌟 NOUVELLE " if created_new else ""
        
        new_name = f"{doc_date}_{predicted_category}_{predicted_sub}_{clean_filename}"
        
        return {
            "original_path": filepath,
            "filename": filename,
            "category": predicted_category,
            "subcategory": predicted_sub,
            "new_name": new_name,
            "status": f"{status_prefix}Classé par IA" if self.api_available else "Classé par nommage",
            "created_new": created_new,
            "reason": reason
        }

# ==================== INTERFACE UTILISATEUR ====================
class TypologyWindow(ctk.CTkToplevel):
    """Fenêtre de gestion de la typologie"""
    def __init__(self, parent, config, on_save):
        super().__init__(parent)
        self.title("Configuration du Plan de Classement")
        self.geometry("800x650")  # Légèrement plus grand pour mieux voir
        self.resizable(True, True)
        
        self.lift()
        self.focus_set()
        self.grab_set()
        
        self.config = config
        self.on_save = on_save
        self.parent_app = parent
        self._build_ui()
        self.draw_items()

    def _build_ui(self):
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(1, weight=1)
        
        # Titre avec bouton de rafraîchissement
        title_frame = ctk.CTkFrame(self, fg_color="transparent")
        title_frame.grid(row=0, column=0, sticky="ew", padx=10, pady=15)
        title_frame.grid_columnconfigure(1, weight=1)
        
        ctk.CTkLabel(title_frame, text="📋 Éditeur de Typologie Métier", 
                    font=("Arial", 20, "bold")).grid(row=0, column=0, sticky="w", padx=5)
        
        # Bouton de rafraîchissement
        ctk.CTkButton(title_frame, text="🔄 Rafraîchir", width=100, height=30,
                     command=self.refresh_display).grid(row=0, column=1, sticky="e", padx=5)
        
        # Info sur le nombre de catégories
        self.stats_label = ctk.CTkLabel(title_frame, text="", 
                                       font=("Arial", 11), text_color="#7f8c8d")
        self.stats_label.grid(row=1, column=0, columnspan=2, sticky="w", padx=5, pady=(5, 0))
        
        # Zone scrollable
        self.scroll = ctk.CTkScrollableFrame(self, width=750, height=450)
        self.scroll.grid(row=1, column=0, sticky="nsew", padx=10, pady=(0, 10))
        self.scroll.grid_columnconfigure(0, weight=1)
        
        # Contrôles
        btn_frame = ctk.CTkFrame(self, fg_color="transparent")
        btn_frame.grid(row=2, column=0, pady=(0, 10))
        
        ctk.CTkButton(btn_frame, text="➕ Ajouter Catégorie", 
                     command=self.add_category, width=150).pack(side="left", padx=5)
        ctk.CTkButton(btn_frame, text="💾 Sauvegarder", fg_color="green",
                     command=self.save_and_close, width=150).pack(side="left", padx=5)
        ctk.CTkButton(btn_frame, text="❌ Fermer", fg_color="gray",
                     command=self.destroy, width=100).pack(side="left", padx=5)

    def refresh_display(self):
        """Rafraîchit l'affichage avec les dernières données"""
        # Recharger la configuration actuelle
        self.config = ConfigManager.load_config()
        self.draw_items()

    def draw_items(self):
        """Affiche la liste des catégories et sous-catégories"""
        # Effacer tout le contenu existant
        for widget in self.scroll.winfo_children():
            widget.destroy()
        
        typology = self.config.get("typology", {})
        
        # Mettre à jour les statistiques
        total_categories = len(typology)
        total_subcategories = sum(len(subs) for subs in typology.values())
        self.stats_label.configure(text=f"📊 {total_categories} catégories • {total_subcategories} sous-catégories")
        
        if not typology:
            # Message si la typologie est vide
            empty_frame = ctk.CTkFrame(self.scroll, height=100)
            empty_frame.pack(fill="x", pady=50)
            ctk.CTkLabel(empty_frame, text="Aucune catégorie définie", 
                        font=("Arial", 16), text_color="#95a5a6").pack(expand=True)
            return
        
        for i, (category, subcategories) in enumerate(typology.items()):
            # Frame de catégorie
            cat_frame = ctk.CTkFrame(self.scroll, corner_radius=8)
            cat_frame.pack(fill="x", pady=8, padx=5)
            cat_frame.grid_columnconfigure(0, weight=1)
            
            # En-tête de catégorie
            cat_header = ctk.CTkFrame(cat_frame, fg_color="transparent")
            cat_header.grid(row=0, column=0, sticky="ew", padx=10, pady=(8, 4))
            cat_header.grid_columnconfigure(0, weight=1)
            
            # Nom de catégorie avec icône
            cat_text = f"📁 {category}"
            if category in ["GENERAL", "AUTRE"] or category not in ["JURIDIQUE", "TECHNIQUE", "COMPTABILITE", "ADMINISTRATIF"]:
                cat_text = f"🌟 {category}"
            
            cat_label = ctk.CTkLabel(cat_header, text=cat_text,
                                   font=("Arial", 16, "bold"), 
                                   text_color="#3498db",
                                   anchor="w")
            cat_label.grid(row=0, column=0, sticky="w", padx=(0, 10))
            
            # Badge du nombre de sous-catégories
            sub_count = len(subcategories)
            count_badge = ctk.CTkLabel(cat_header, text=f"{sub_count} sous-cat.",
                                     font=("Arial", 10), 
                                     text_color="#7f8c8d",
                                     fg_color="#2c3e50",
                                     corner_radius=10)
            count_badge.grid(row=0, column=1, padx=5)
            
            # Boutons catégorie
            btn_frame = ctk.CTkFrame(cat_header, fg_color="transparent")
            btn_frame.grid(row=0, column=2, padx=5)
            
            ctk.CTkButton(btn_frame, text="➕", width=35, height=30,
                         command=lambda c=category: self.add_subcategory(c)).pack(side="left", padx=2)
            ctk.CTkButton(btn_frame, text="✏️", width=35, height=30, fg_color="#f39c12",
                         command=lambda c=category: self.edit_category(c)).pack(side="left", padx=2)
            ctk.CTkButton(btn_frame, text="❌", width=35, height=30, fg_color="#e74c3c",
                         command=lambda c=category: self.delete_category(c)).pack(side="left", padx=2)
            
            # Sous-catégories
            if subcategories:
                sub_frame = ctk.CTkFrame(cat_frame, fg_color="transparent")
                sub_frame.grid(row=1, column=0, sticky="ew", padx=20, pady=(0, 8))
                
                for j, sub in enumerate(subcategories):
                    sub_row = ctk.CTkFrame(sub_frame, fg_color="transparent", height=35)
                    sub_row.pack(fill="x", pady=2)
                    sub_row.grid_columnconfigure(0, weight=1)
                    
                    ctk.CTkLabel(sub_row, text="  └ 📄", 
                               text_color="#95a5a6", font=("Arial", 12)).grid(row=0, column=0, sticky="w", padx=(0, 5))
                    
                    sub_label = ctk.CTkLabel(sub_row, text=sub, 
                                           font=("Arial", 13), anchor="w")
                    sub_label.grid(row=0, column=1, sticky="w", padx=5)
                    
                    # Boutons sous-catégorie
                    sub_btn_frame = ctk.CTkFrame(sub_row, fg_color="transparent")
                    sub_btn_frame.grid(row=0, column=2, sticky="e")
                    
                    ctk.CTkButton(sub_btn_frame, text="✏️", width=30, height=26,
                                command=lambda c=category, s=sub: self.edit_subcategory(c, s)).pack(side="left", padx=2)
                    ctk.CTkButton(sub_btn_frame, text="❌", width=30, height=26, fg_color="#e74c3c",
                                command=lambda c=category, s=sub: self.delete_subcategory(c, s)).pack(side="left", padx=2)
            else:
                empty_sub_frame = ctk.CTkFrame(cat_frame, fg_color="transparent")
                empty_sub_frame.grid(row=1, column=0, sticky="ew", padx=30, pady=(0, 8))
                
                ctk.CTkLabel(empty_sub_frame, text="  └ (Aucune sous-catégorie)", 
                           text_color="#bdc3c7", font=("Arial", 11, "italic")).pack(anchor="w")

    def add_category(self):
        name = simpledialog.askstring("Nouvelle Catégorie", 
                                    "Nom de la catégorie (ex: URBANISME) :", 
                                    parent=self)
        if name:
            name = name.upper().strip()
            if name and name not in self.config.get("typology", {}):
                self.config.setdefault("typology", {})[name] = []
                self.draw_items()
                messagebox.showinfo("Succès", f"Catégorie '{name}' ajoutée avec succès!", parent=self)

    def edit_category(self, old_name):
        new_name = simpledialog.askstring("Modifier Catégorie", 
                                        f"Nouveau nom pour '{old_name}' :", 
                                        parent=self)
        if new_name:
            new_name = new_name.upper().strip()
            if new_name and new_name != old_name:
                typology = self.config.get("typology", {})
                if old_name in typology:
                    typology[new_name] = typology.pop(old_name)
                    self.draw_items()
                    messagebox.showinfo("Succès", f"Catégorie renommée: '{old_name}' → '{new_name}'", parent=self)

    def delete_category(self, category):
        if messagebox.askyesno("Confirmation", 
                             f"Supprimer la catégorie '{category}' et toutes ses sous-catégories ?", 
                             parent=self):
            typology = self.config.get("typology", {})
            if category in typology:
                del typology[category]
                self.draw_items()
                messagebox.showinfo("Succès", f"Catégorie '{category}' supprimée", parent=self)

    def add_subcategory(self, category):
        name = simpledialog.askstring("Nouvelle Sous-Catégorie", 
                                    f"Sous-catégorie pour '{category}' :", 
                                    parent=self)
        if name:
            name = name.strip()
            typology = self.config.get("typology", {})
            if category in typology and name not in typology[category]:
                typology[category].append(name)
                self.draw_items()
                messagebox.showinfo("Succès", f"Sous-catégorie '{name}' ajoutée à '{category}'", parent=self)

    def edit_subcategory(self, category, old_sub):
        new_sub = simpledialog.askstring("Modifier Sous-Catégorie", 
                                       f"Nouveau nom pour '{old_sub}' :", 
                                       parent=self)
        if new_sub:
            new_sub = new_sub.strip()
            typology = self.config.get("typology", {})
            if category in typology and old_sub in typology[category]:
                idx = typology[category].index(old_sub)
                typology[category][idx] = new_sub
                self.draw_items()
                messagebox.showinfo("Succès", f"Sous-catégorie renommée: '{old_sub}' → '{new_sub}'", parent=self)

    def delete_subcategory(self, category, subcategory):
        if messagebox.askyesno("Confirmation", 
                             f"Supprimer la sous-catégorie '{subcategory}' ?", 
                             parent=self):
            typology = self.config.get("typology", {})
            if category in typology and subcategory in typology[category]:
                typology[category].remove(subcategory)
                self.draw_items()
                messagebox.showinfo("Succès", f"Sous-catégorie '{subcategory}' supprimée", parent=self)

    def save_and_close(self):
        ConfigManager.save_config(self.config)
        if self.on_save:
            self.on_save()
        messagebox.showinfo("Sauvegarde", "Plan de classement sauvegardé avec succès!", parent=self)
        self.destroy()

class MainApp(ctk.CTk):
    """Application principale"""
    def __init__(self):
        super().__init__()
        
        self.config = ConfigManager.load_config()
        self.file_index = ConfigManager.load_index()
        self.classification_engine = ClassificationEngine()
        self.current_files = []
        self.new_categories_created = []  # Pour suivre les nouvelles catégories
        self.typology_window = None  # Référence à la fenêtre de typologie
        
        self._setup_appearance()
        self._setup_ui()
        self._update_stats()

    def _setup_appearance(self):
        ctk.set_appearance_mode("dark")
        ctk.set_default_color_theme("blue")
        
        self.title("MALKOGED AI - TVD")
        self.geometry("1200x800")
        
        # 1. On cache la fenêtre pendant les calculs pour éviter le "clignotement"
        self.withdraw() 
        
        # 2. Définition d'une taille par défaut généreuse (75% de l'écran)
        screen_width = self.winfo_screenwidth()
        screen_height = self.winfo_screenheight()
        
        width = int(screen_width * 0.8)
        height = int(screen_height * 0.8)
        
        # 3. Centrage précis
        x = (screen_width // 2) - (width // 2)
        y = (screen_height // 2) - (height // 2)
        
        self.geometry(f'{width}x{height}+{x}+{y}')
        
        # 5. On réaffiche la fenêtre une fois prête
        self.after(200, self.deiconify)

    def _setup_ui(self):
        # Configuration de la grille
        self.grid_columnconfigure(1, weight=1)
        self.grid_rowconfigure(0, weight=1)
        
        # ============ SIDEBAR ============
        self.sidebar = ctk.CTkFrame(self, width=280, corner_radius=0)
        self.sidebar.grid(row=0, column=0, sticky="nsew", padx=5, pady=5)
        
        # Logo / Titre
        ctk.CTkLabel(self.sidebar, text="MALKOGED AI", 
                    font=("Arial", 24, "bold")).pack(pady=(30, 5))
        ctk.CTkLabel(self.sidebar, text="IMMO-MOUSQUETAIRES_GED", 
                    font=("Arial", 12), text_color="#7f8c8d").pack(pady=(0, 20))
        
        # ============ NOUVEAU : BOUTONS D'IMPORTATION ============
        import_frame = ctk.CTkFrame(self.sidebar, fg_color="transparent")
        import_frame.pack(pady=10, padx=20, fill="x")
        
        ctk.CTkButton(import_frame, text="📄 Importer Fichiers", 
                      command=self.import_files, height=35).pack(pady=5, fill="x")
        
        ctk.CTkButton(import_frame, text="📁 Importer Dossier", 
                      command=self.import_folder, height=35).pack(pady=5, fill="x")
        
        # ============ BOUTONS D'ACTION ============
        action_frame = ctk.CTkFrame(self.sidebar, fg_color="transparent")
        action_frame.pack(pady=10, padx=20, fill="x")
        
        ctk.CTkButton(action_frame, text="🚀 Traiter & Classer", 
                     command=self.process_imported, height=40).pack(pady=10, fill="x")
        ctk.CTkButton(action_frame, text="🔍 Vérifier Doublons", 
                     command=self.check_duplicates, height=40, fg_color="#f39c12").pack(pady=10, fill="x")
        
        # Options
        options_frame = ctk.CTkFrame(self.sidebar, fg_color="transparent")
        options_frame.pack(pady=20, padx=20, fill="x")
        
        self.auto_delete_var = ctk.BooleanVar(value=self.config.get("auto_delete", False))
        self.auto_delete_check = ctk.CTkCheckBox(options_frame, text="Supprimer après archivage",
                                               variable=self.auto_delete_var)
        self.auto_delete_check.pack(anchor="w", pady=5)
        
        self.api_active_var = ctk.BooleanVar(value=self.config.get("api_active", True))
        self.api_check = ctk.CTkCheckBox(options_frame, text="Activer IA DeepSeek",
                                       variable=self.api_active_var,
                                       command=self.toggle_api)
        self.api_check.pack(anchor="w", pady=5)
        
        # NOUVELLE OPTION : Création automatique de catégories
        self.auto_create_var = ctk.BooleanVar(value=self.config.get("auto_create_categories", True))
        self.auto_create_check = ctk.CTkCheckBox(options_frame, text="Créer catégories auto",
                                               variable=self.auto_create_var,
                                               command=self.toggle_auto_create)
        self.auto_create_check.pack(anchor="w", pady=5)
        
        # Configuration
        config_frame = ctk.CTkFrame(self.sidebar, fg_color="transparent")
        config_frame.pack(pady=10, padx=20, fill="x")
        
        ctk.CTkButton(config_frame, text="⚙️ Plan de Classement", 
                     command=self.open_typology, fg_color="#34495e", height=35).pack(pady=5, fill="x")
        ctk.CTkButton(config_frame, text="🔌 Tester API", 
                     command=self.test_api, fg_color="#27ae60", height=35).pack(pady=5, fill="x")
        
        # Statistiques
        stats_frame = ctk.CTkFrame(self.sidebar)
        stats_frame.pack(side="bottom", fill="x", padx=20, pady=20)
        
        self.stats_label = ctk.CTkLabel(stats_frame, text="", font=("Arial", 11))
        self.stats_label.pack(pady=10)
        
        # ============ MAIN AREA ============
        main_container = ctk.CTkFrame(self, fg_color="transparent")
        main_container.grid(row=0, column=1, sticky="nsew", padx=10, pady=10)
        main_container.grid_columnconfigure(0, weight=1)
        main_container.grid_rowconfigure(1, weight=1)
        
        # En-tête
        header = ctk.CTkFrame(main_container, height=60)
        header.grid(row=0, column=0, sticky="ew", pady=(0, 10))
        header.grid_columnconfigure(1, weight=1)
        
        ctk.CTkLabel(header, text="Journal de Traitement", 
                    font=("Arial", 22, "bold")).grid(row=0, column=0, sticky="w", padx=10)
        
        self.status_label = ctk.CTkLabel(header, text="Prêt", 
                                       font=("Arial", 12))
        self.status_label.grid(row=0, column=1, sticky="e", padx=10)
        
        # Tableau des résultats
        self._create_results_table(main_container)

    def _create_results_table(self, parent):
        """Crée le tableau des résultats"""
        table_frame = ctk.CTkFrame(parent)
        table_frame.grid(row=1, column=0, sticky="nsew")
        table_frame.grid_columnconfigure(0, weight=1)
        table_frame.grid_rowconfigure(1, weight=1)
        
        # En-têtes
        headers = ["Fichier", "Catégorie", "Sous-Catégorie", "Statut", "Actions"]
        header_frame = ctk.CTkFrame(table_frame, height=40)
        header_frame.grid(row=0, column=0, sticky="ew", padx=5, pady=(0, 5))
        
        for i, header in enumerate(headers):
            width = 200 if i == 0 else (150 if i < 4 else 100)
            ctk.CTkLabel(header_frame, text=header, font=("Arial", 13, "bold"),
                        width=width).grid(row=0, column=i, padx=2)
        
        # Zone scrollable
        self.results_scroll = ctk.CTkScrollableFrame(table_frame, height=550)
        self.results_scroll.grid(row=1, column=0, sticky="nsew", padx=5)
        
        # Configuration colonnes
        for i in range(5):
            self.results_scroll.grid_columnconfigure(i, weight=1)

    def _update_stats(self):
        """Met à jour les statistiques affichées"""
        total_files = len(self.file_index)
        typology = self.config.get("typology", {})
        typology_size = len(typology)
        total_subcategories = sum(len(subs) for subs in typology.values())
        
        stats_text = f"📊 Statistiques\n"
        stats_text += f"Fichiers indexés: {total_files}\n"
        stats_text += f"Catégories: {typology_size}\n"
        stats_text += f"Sous-catégories: {total_subcategories}\n"
        stats_text += f"API: {'✅ Active' if self.config.get('api_active', True) else '❌ Inactive'}\n"
        stats_text += f"Auto-création: {'✅ ON' if self.config.get('auto_create_categories', True) else '❌ OFF'}"
        
        self.stats_label.configure(text=stats_text)

    def toggle_api(self):
        """Active/désactive l'API"""
        self.config["api_active"] = self.api_active_var.get()
        ConfigManager.save_config(self.config)
        self.classification_engine.reload_typology()
        self._update_stats()

    def toggle_auto_create(self):
        """Active/désactive la création automatique de catégories"""
        self.config["auto_create_categories"] = self.auto_create_var.get()
        ConfigManager.save_config(self.config)
        self.classification_engine.reload_typology()
        self._update_stats()

    def open_typology(self):
        """Ouvre ou rafraîchit la fenêtre de configuration de typologie"""
        # Si la fenêtre existe déjà, la détruire d'abord
        if hasattr(self, 'typology_window') and self.typology_window is not None:
            try:
                self.typology_window.destroy()
            except:
                pass
        
        # Créer une nouvelle fenêtre avec les données à jour
        self.typology_window = TypologyWindow(self, self.config, self._on_typology_saved)
        # Mettre à jour l'affichage immédiatement
        self.typology_window.refresh_display()

    def _on_typology_saved(self):
        """Callback après sauvegarde de la typologie"""
        # Recharger la configuration
        self.config = ConfigManager.load_config()
        # Recharger la typologie dans le moteur de classification
        self.classification_engine.reload_typology()
        # Mettre à jour les stats
        self._update_stats()
        
        # Si la fenêtre de typologie est ouverte, la rafraîchir
        if hasattr(self, 'typology_window') and self.typology_window is not None:
            try:
                self.typology_window.refresh_display()
            except:
                pass

    def refresh_typology_window(self):
        """Rafraîchit la fenêtre de typologie si elle est ouverte"""
        if hasattr(self, 'typology_window') and self.typology_window is not None:
            try:
                self.typology_window.refresh_display()
            except:
                pass

    def test_api(self):
        """Teste la connexion à l'API DeepSeek"""
        if not self.config.get("api_active", True):
            messagebox.showwarning("API", "L'API est désactivée dans les options.")
            return
            
        try:
            headers = {
                "Authorization": f"Bearer {API_KEY}",
                "Content-Type": "application/json"
            }
            
            payload = {
                "model": "deepseek-chat",
                "messages": [{"role": "user", "content": "Réponds par 'API OK'"}],
                "max_tokens": 10
            }
            
            response = requests.post(DEEPSEEK_API_URL, headers=headers, json=payload, timeout=10)
            response.raise_for_status()
            
            result = response.json()
            messagebox.showinfo("API Test", f"✅ Connexion API réussie!\nModèle: {result.get('model', 'Inconnu')}")
            
        except Exception as e:
            messagebox.showerror("API Error", f"❌ Échec connexion API:\n{str(e)}")

    def import_files(self):
        """Sélecteur mis à jour pour inclure Office"""
        supported_ext = [
            ("Tous les documents", "*.pdf *.docx *.xlsx *.pptx *.mp3 *.wav *.mp4 *.mov *.avi *.jpg *.jpeg *.png"),
            ("Documents Office", "*.docx *.xlsx *.pptx"),
            ("PDF", "*.pdf"),
            ("Multimédia", "*.mp3 *.wav *.mp4 *.mov *.avi"),
            ("Images", "*.jpg *.jpeg *.png"),
            ("Tous les fichiers", "*.*")
        ]
        files = filedialog.askopenfilenames(title="Sélectionnez les documents", filetypes=supported_ext)
        if files:
            self.start_processing(list(files))

    def import_folder(self):
        """Sélection d'un dossier complet"""
        folder = filedialog.askdirectory(title="Sélectionnez le dossier source")
        if folder:
            # Filtre avec tous les formats supportés
            supported_ext = ('.pdf', '.docx', '.xlsx', '.pptx', '.mp3', '.wav', 
                           '.mp4', '.mov', '.avi', '.jpg', '.jpeg', '.png')
            files = []
            for root, _, filenames in os.walk(folder):
                for filename in filenames:
                    if filename.lower().endswith(supported_ext):
                        files.append(os.path.join(root, filename))
            
            if files:
                self.start_processing(files)
            else:
                messagebox.showwarning("Aucun fichier", "Aucun fichier compatible trouvé dans ce dossier.")

    def process_imported(self):
        """Traiter les fichiers déjà importés ou sélectionner de nouveaux"""
        if not hasattr(self, 'current_files') or not self.current_files:
            messagebox.showinfo("Aucun fichier", 
                              "Aucun fichier à traiter. Veuillez d'abord importer des fichiers.")
            return
        
        dest_dir = filedialog.askdirectory(title="Sélectionnez le dossier de destination (Archives)")
        if not dest_dir:
            return

        # Mise à jour de la config
        self.config["last_destination"] = dest_dir
        ConfigManager.save_config(self.config)

        # Réinitialiser la liste des nouvelles catégories
        self.new_categories_created = []
        
        # Lancement du thread
        threading.Thread(target=self._process_files_thread, 
                         args=(self.current_files, dest_dir), daemon=True).start()

    def start_processing(self, file_list):
        """Lance le traitement commun pour fichiers ou dossiers"""
        dest_dir = filedialog.askdirectory(title="Sélectionnez le dossier de destination (Archives)")
        if not dest_dir:
            return

        # Mise à jour de la config
        self.config["last_destination"] = dest_dir
        ConfigManager.save_config(self.config)

        # Nettoyage interface
        self.after(0, self._clear_results)
        self.current_files = file_list
        self.new_categories_created = []
        
        # Lancement du thread
        threading.Thread(target=self._process_files_thread, 
                         args=(file_list, dest_dir), daemon=True).start()

    def _process_files_thread(self, file_list, dest_dir):
        """Thread de traitement des fichiers"""
        if not file_list:
            self.after(0, lambda: messagebox.showwarning("Aucun fichier", 
                                                       "Aucun fichier à traiter."))
            return
        
        # Fenêtre de progression
        self.after(0, self._show_progress, len(file_list))
        
        # Traitement
        processed = 0
        duplicates = 0
        errors = 0
        
        for i, filepath in enumerate(file_list):
            # Mise à jour progression
            self.after(0, self._update_progress, i + 1, len(file_list))
            
            try:
                result = self._process_single_file(filepath, dest_dir)
                
                if result["status"] == "DOUBLON":
                    duplicates += 1
                elif "ERREUR" in result["status"]:
                    errors += 1
                else:
                    processed += 1
                    # Si une nouvelle catégorie a été créée, la suivre
                    if result.get("created_new", False):
                        new_cat = {
                            "category": result["category"],
                            "subcategory": result["subcategory"],
                            "file": result["filename"],
                            "reason": result.get("reason", "")
                        }
                        self.new_categories_created.append(new_cat)
                
                self.after(0, self._add_result_row, result)
                
            except Exception as e:
                print(f"Erreur traitement {filepath}: {e}")
                errors += 1
        
        # Fermeture progression
        self.after(0, self._hide_progress)
        
        # Sauvegarde index
        ConfigManager.save_index(self.file_index)
        
        # Rafraîchir la configuration pour avoir les dernières catégories
        self.config = ConfigManager.load_config()
        
        # Rafraîchir la fenêtre de typologie si elle est ouverte
        self.after(0, self.refresh_typology_window)
        
        # Affichage résultats avec nouvelles catégories
        self.after(0, lambda: self._show_results(processed, duplicates, errors, file_list))

    def _process_single_file(self, filepath, dest_dir):
        """Traite un fichier individuel"""
        filename = os.path.basename(filepath)
        
        # 1. Vérification doublon
        file_hash = DuplicateManager.get_file_hash(filepath)
        is_duplicate = DuplicateManager.is_duplicate(file_hash, self.file_index)
        
        if is_duplicate:
            return {
                "filename": filename,
                "category": "DOUBLON",
                "subcategory": "",
                "status": f"DOUBLON ({os.path.basename(self.file_index[file_hash])[:20]}...)",
                "color": "orange",
                "path": filepath,
                "is_duplicate": True,
                "created_new": False
            }
        
        # 2. Classification avec création automatique
        classification = self.classification_engine.analyze_document(filepath)
        
        # 3. Préparation destination
        final_dir = os.path.join(dest_dir, classification["category"], classification["subcategory"])
        os.makedirs(final_dir, exist_ok=True)
        
        dest_path = os.path.join(final_dir, classification["new_name"])
        
        # 4. Copie
        shutil.copy2(filepath, dest_path)
        
        # 5. Vérification intégrité
        dest_hash = DuplicateManager.get_file_hash(dest_path)
        if dest_hash != file_hash:
            return {
                "filename": filename,
                "category": classification["category"],
                "subcategory": classification["subcategory"],
                "status": "ERREUR Intégrité",
                "color": "red",
                "path": dest_path,
                "is_duplicate": False,
                "created_new": False
            }
        
        # 6. Mise à jour index
        self.file_index[file_hash] = dest_path
        
        # 7. Tagging métadonnées (si fichier audio/vidéo)
        if dest_path.lower().endswith(('.mp3', '.mp4', '.m4a')):
            try:
                MetadataManager.tag_file(dest_path, classification["category"], classification["subcategory"])
            except:
                pass
        
        # 8. Suppression source si option activée
        source_deleted = False
        if self.auto_delete_var.get():
            try:
                os.remove(filepath)
                source_deleted = True
            except:
                pass
        
        # 9. Retour résultat
        status = f"{classification['status']}{' (Source supprimée)' if source_deleted else ''}"
        
        return {
            "filename": filename,
            "category": classification["category"],
            "subcategory": classification["subcategory"],
            "status": status,
            "color": "#27ae60" if not classification.get("created_new", False) else "#f39c12",
            "path": dest_path,
            "is_duplicate": False,
            "created_new": classification.get("created_new", False),
            "reason": classification.get("reason", ""),
            "new_name": classification["new_name"]
        }

    def check_duplicates(self):
        """Vérifie les doublons dans un dossier"""
        source_dir = filedialog.askdirectory(title="Sélectionnez le dossier à vérifier")
        if not source_dir:
            return
            
        # Recherche doublons
        duplicates_found = []
        file_hashes = {}
        
        for root, _, files in os.walk(source_dir):
            for file in files:
                filepath = os.path.join(root, file)
                file_hash = DuplicateManager.get_file_hash(filepath)
                
                if file_hash:
                    if file_hash in file_hashes:
                        duplicates_found.append((file, file_hashes[file_hash]))
                    else:
                        file_hashes[file_hash] = file
        
        if duplicates_found:
            message = f"Doublons trouvés: {len(duplicates_found)}\n\n"
            for dup in duplicates_found[:10]:  # Limite à 10 affichages
                message += f"- {dup[0]} (identique à {dup[1]})\n"
            
            if len(duplicates_found) > 10:
                message += f"\n... et {len(duplicates_found) - 10} autres"
            
            messagebox.showwarning("Doublons détectés", message)
        else:
            messagebox.showinfo("Vérification", "Aucun doublon détecté.")

    def _clear_results(self):
        """Vide le tableau des résultats"""
        for widget in self.results_scroll.winfo_children():
            widget.destroy()

    def _add_result_row(self, result):
        """Ajoute une ligne au tableau des résultats"""
        row_idx = len(self.results_scroll.winfo_children()) // 5
        
        # Fichier
        ctk.CTkLabel(self.results_scroll, text=result["filename"][:30], 
                    anchor="w").grid(row=row_idx, column=0, padx=2, pady=2, sticky="w")
        
        # Catégorie avec badge si nouvelle
        cat_text = result["category"]
        if result.get("created_new", False):
            cat_text = f"🌟 {cat_text}"
        
        cat_color = result["color"]
        cat_label = ctk.CTkLabel(self.results_scroll, text=cat_text, 
                                anchor="w", text_color=cat_color)
        cat_label.grid(row=row_idx, column=1, padx=2, pady=2, sticky="w")
        
        # Sous-catégorie
        ctk.CTkLabel(self.results_scroll, text=result["subcategory"], 
                    anchor="w").grid(row=row_idx, column=2, padx=2, pady=2, sticky="w")
        
        # Statut
        ctk.CTkLabel(self.results_scroll, text=result["status"], 
                    anchor="w").grid(row=row_idx, column=3, padx=2, pady=2, sticky="w")
        
        # Actions
        if not result.get("is_duplicate", True) and "path" in result:
            btn_frame = ctk.CTkFrame(self.results_scroll, fg_color="transparent")
            btn_frame.grid(row=row_idx, column=4, padx=2, pady=2)
            
            ctk.CTkButton(btn_frame, text="📂", width=30,
                         command=lambda p=result["path"]: os.startfile(os.path.dirname(p))).pack(side="left", padx=2)

    def _show_progress(self, total_files):
        """Affiche la fenêtre de progression"""
        self.progress_window = ctk.CTkToplevel(self)
        self.progress_window.title("Traitement en cours")
        self.progress_window.geometry("400x150")
        self.progress_window.transient(self)
        self.progress_window.grab_set()
        
        ctk.CTkLabel(self.progress_window, text="Analyse et archivage...", 
                    font=("Arial", 14)).pack(pady=20)
        
        self.progress_bar = ctk.CTkProgressBar(self.progress_window, width=350)
        self.progress_bar.pack(pady=10)
        self.progress_bar.set(0)
        
        self.progress_label = ctk.CTkLabel(self.progress_window, text=f"0/{total_files}")
        self.progress_label.pack()

    def _update_progress(self, current, total):
        """Met à jour la barre de progression"""
        if hasattr(self, 'progress_bar'):
            progress = current / total
            self.progress_bar.set(progress)
            self.progress_label.configure(text=f"{current}/{total}")
            self.progress_window.update()

    def _hide_progress(self):
        """Cache la fenêtre de progression"""
        if hasattr(self, 'progress_window'):
            self.progress_window.destroy()
            del self.progress_window

    def _show_results(self, processed, duplicates, errors, file_list):
        """Affiche le résumé du traitement avec nouvelles catégories"""
        message = f"Traitement terminé !\n\n"
        message += f"✅ Fichiers traités: {processed}\n"
        message += f"🔄 Doublons ignorés: {duplicates}\n"
        message += f"❌ Erreurs: {errors}\n\n"
        
        if self.auto_delete_var.get():
            message += "⚠️ Les fichiers sources ont été supprimés.\n\n"
        
        # Afficher les nouvelles catégories créées
        if self.new_categories_created:
            message += "🌟 NOUVELLES CATÉGORIES CRÉÉES :\n"
            for i, new_cat in enumerate(self.new_categories_created[:5]):  # Limite à 5
                message += f"{i+1}. {new_cat['category']} > {new_cat['subcategory']}\n"
                if new_cat.get('reason'):
                    message += f"   Raison: {new_cat['reason'][:50]}...\n"
            
            if len(self.new_categories_created) > 5:
                message += f"... et {len(self.new_categories_created) - 5} autres\n"
            
            message += f"\nTotal: {len(self.new_categories_created)} nouvelle(s) catégorie(s)\n"
        
        # Proposer de sauvegarder le nouveau plan de classement
        if self.new_categories_created:
            message += "\n📋 Le plan de classement a été automatiquement mis à jour."
            message += "\n\nCliquez sur 'Plan de Classement' pour voir les nouvelles catégories."
        
        messagebox.showinfo("Résultats", message)
        self.status_label.configure(text=f"Terminé - {processed} fichiers traités")
        self._update_stats()  # Met à jour les stats avec les nouvelles catégories

# ==================== LANCEMENT ====================
if __name__ == "__main__":
    # Installation requise :
    # pip install customtkinter pdfplumber requests mutagen pillow python-docx openpyxl python-pptx
    
    app = MainApp()
    app.mainloop()
